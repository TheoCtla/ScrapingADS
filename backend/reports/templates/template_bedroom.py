"""
Template Bedroom — Hérite de TemplateEmma.
Seule différence : slide Conversions dédiée (Appels + Itinéraires)
entre Meta Ads et Synthèse.
"""

from pathlib import Path
from typing import Optional

from pptx import Presentation

from backend.reports.templates.template_emma import TemplateEmma
from backend.reports.styles import (
    PALETTE, format_number, format_currency,
)

ASSETS_DIR = Path(__file__).parent.parent / "assets"

COVER_IMAGES = [
    ("bedroom", "img_bedroom.jpg"),
]


def _resolve_cover_image(client_name: str) -> Optional[str]:
    lower_name = client_name.lower()
    for keyword, filename in COVER_IMAGES:
        if keyword in lower_name:
            path = ASSETS_DIR / filename
            if path.exists():
                return str(path)
    return None


class TemplateBedroom(TemplateEmma):
    """Template pour les magasins Bedroom."""

    def generate(self, data: dict) -> Presentation:
        google = data.get("google_ads", {})
        meta = data.get("meta_ads", {})
        conv = data.get("conversions", {})
        general = data.get("general", {})
        history = data.get("history", [])

        g_curr = google.get("current", {})
        g_prev = google.get("previous", {})
        m_curr = meta.get("current", {})
        m_prev = meta.get("previous", {})
        c_curr = conv.get("current", {})
        c_prev = conv.get("previous", {})
        gen_curr = general.get("current", {})
        gen_prev = general.get("previous", {})

        client = data.get("client", "")
        month_fr = data.get("month_fr", "")

        has_google = self._safe_get(g_curr, "Cout Google ADS") > 0
        has_meta = self._safe_get(m_curr, "Cout Facebook ADS") > 0
        has_history = history and len(history) >= 2
        has_itineraires = self._get_iti(c_curr) > 0
        has_appels = self._safe_get(c_curr, "Appels Téléphoniques") > 0

        # Slide 1 — Titre
        cover_image = data.get("cover_image") or _resolve_cover_image(client)
        self._title_slide(client, month_fr, image_path=cover_image)

        # Slide 2 — Récap (Diffusion, Google cost, Meta cost)
        self._slide_recap(gen_curr, gen_prev, g_curr, g_prev, m_curr, m_prev, month_fr)

        # Slide 3 — Google Ads
        if has_google:
            self._slide_google(g_curr, g_prev, month_fr)
            if has_history:
                self._slide_evo_quad(history, [
                    ("Cout Google ADS", "Coût"),
                    ("Total Impressions", "Impressions"),
                    ("Total Clic", "Clics totaux"),
                    ("Total CPC moyen", "CPC Moyen"),
                ], "Google Ads")

        # Slide 4 — Meta Ads
        if has_meta:
            self._slide_meta(m_curr, m_prev, month_fr)
            if has_history:
                self._slide_evo_quad(history, [
                    ("Cout Facebook ADS", "Coût"),
                    ("Impressions Meta", "Impressions"),
                    ("Clics Meta", "Clics"),
                    ("CPC Meta", "CPC"),
                ], "Meta Ads - Facebook et Instagram")

        # Slide 5 — Conversions (Appels + Itinéraires)
        if has_appels or has_itineraires:
            self._slide_conversions(c_curr, c_prev, month_fr, has_appels, has_itineraires)
            if has_history and (has_itineraires or has_appels):
                self._slide_evo_conversions(history, has_itineraires, has_appels)

        # Slide 6 — Synthèse
        self._slide_synthese(gen_curr, gen_prev, c_curr, c_prev, month_fr,
                             has_itineraires, has_appels)

        # Slide 7 — Fin
        self._end_slide()

        return self.prs

    # ──────────────────────────────────────────
    # Slide Conversions (spécifique Bedroom)
    # ──────────────────────────────────────────

    def _slide_conversions(self, c_curr, c_prev, month_fr, has_appels, has_itineraires):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Conversions", month_fr, accent_color=PALETTE["gold"])

        row = []
        if has_appels:
            appels = self._safe_get(c_curr, "Appels Téléphoniques")
            appels_p = self._safe_get(c_prev, "Appels Téléphoniques")
            row.append({"label": "Appels téléphoniques", "value": format_number(appels),
                        "current": appels, "previous": appels_p,
                        "sub_label": "Nombre de clics trackés sur les boutons contacts",
                        "accent": PALETTE["gold"]})
        if has_itineraires:
            iti = self._get_iti(c_curr)
            iti_p = self._get_iti(c_prev)
            row.append({"label": "Itinéraires", "value": format_number(iti),
                        "current": iti, "previous": iti_p,
                        "sub_label": "Nombre de clics trackés sur les boutons itinéraires",
                        "accent": PALETTE["gold"]})

        self._hero_row(slide, row, y=2.5, n_cols=len(row))

    # ──────────────────────────────────────────
    # Synthèse Bedroom (sans coût/conversion)
    # ──────────────────────────────────────────

    def _slide_synthese(self, gen_curr, gen_prev, c_curr, c_prev, month_fr,
                        has_itineraires, has_appels):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Synthèse", month_fr)

        cout_all = self._safe_get(gen_curr, "COUT ALL")
        cout_all_p = self._safe_get(gen_prev, "COUT ALL")

        row1 = [
            {"label": "Achat Média Total", "value": format_currency(cout_all),
             "current": cout_all, "previous": cout_all_p, "tooltip": "COUT ALL",
             "accent": PALETTE["gold"]},
        ]
        self._hero_row(slide, row1, y=1.3, n_cols=1)

        row2 = []
        if has_appels:
            appels = self._safe_get(c_curr, "Appels Téléphoniques")
            appels_p = self._safe_get(c_prev, "Appels Téléphoniques")
            row2.append({"label": "Appels téléphoniques", "value": format_number(appels),
                         "current": appels, "previous": appels_p,
                         "sub_label": "Nombre de clics trackés sur les boutons contacts",
                         "accent": PALETTE["gold"]})
        if has_itineraires:
            iti = self._get_iti(c_curr)
            iti_p = self._get_iti(c_prev)
            row2.append({"label": "Itinéraires", "value": format_number(iti),
                         "current": iti, "previous": iti_p,
                         "sub_label": "Nombre de clics trackés sur les boutons itinéraires",
                         "accent": PALETTE["gold"]})

        if row2:
            self._hero_row(slide, row2, y=4.2, n_cols=len(row2))
