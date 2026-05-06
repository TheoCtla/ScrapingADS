"""
Template LyleOO — Même charte graphique que les autres templates.
Spécificités : centré sur Meta Ads, slide interactions manuelles,
pas de slide Google détaillée, pas de Microsoft.
"""

import os
import tempfile
import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from backend.reports.templates.base import BaseTemplate
from backend.reports.styles import (
    PALETTE, FUNCTIONAL_COLORS, MARGIN, SLIDE_WIDTH, SLIDE_HEIGHT,
    METRIC_TOOLTIPS, INVERSE_METRICS,
    hex_to_rgb, format_currency, format_number, format_percentage, calc_variation,
)
from backend.reports import charts

ASSETS_DIR = Path(__file__).parent.parent / "assets"
LOGO_PATH = ASSETS_DIR / "tarmaac_logo.png"

COVER_IMAGE = "img_lyleoo.jpg"

# Interactions par campagne (placeholders manuels)
INTEREST_ROWS = [
    "Interaction Facebook",
    "Interaction Insta",
    "Budget",
]

RETARGET_ROWS = [
    "Interaction Facebook",
    "Interaction Insta",
    "Budget",
]


class TemplateLyleoo(BaseTemplate):
    """Template pour LyleOO — centré Meta Ads."""

    @staticmethod
    def _has_data(section: dict) -> bool:
        return any(
            v not in (0, None, "", "0", "0.0")
            for k, v in section.items()
            if k not in ("month", "month_fr")
        )

    def generate(self, data: dict) -> Presentation:
        google = data.get("google_ads", {})
        meta = data.get("meta_ads", {})
        general = data.get("general", {})
        history = data.get("history", [])

        g_curr = google.get("current", {})
        g_prev = google.get("previous", {})
        m_curr = meta.get("current", {})
        m_prev = meta.get("previous", {})
        gen_curr = general.get("current", {})
        gen_prev = general.get("previous", {})

        client = data.get("client", "")
        month_fr = data.get("month_fr", "")
        prev_month_fr = data.get("previous_month_fr", "")

        has_meta = self._safe_get(m_curr, "Cout Facebook ADS") > 0
        has_history = history and len(history) >= 2

        # Slide 1 — Titre
        cover_image = data.get("cover_image")
        if not cover_image:
            path = ASSETS_DIR / COVER_IMAGE
            if path.exists():
                cover_image = str(path)
        self._title_slide(client, month_fr, image_path=cover_image)

        # Slide 2 — Récap
        self._slide_recap(gen_curr, gen_prev, g_curr, g_prev, m_curr, m_prev, month_fr)

        # Slide 3 — Meta Ads détails
        if has_meta:
            self._slide_meta(m_curr, m_prev, month_fr)
            if has_history:
                self._slide_evo_quad(history, [
                    ("Cout Facebook ADS", "Coût"),
                    ("Impressions Meta", "Impressions"),
                    ("Clics Meta", "Clics"),
                    ("CPC Meta", "CPC"),
                ], "Meta Ads - Facebook et Instagram")

        # Slide 4 — Interactions par campagne (placeholders)
        self._slide_interactions(month_fr, prev_month_fr)

        # Slide 5 — Synthèse
        self._slide_synthese(m_curr, m_prev, month_fr)

        # Slide 6 — Fin
        self._end_slide()

        return self.prs

    # ──────────────────────────────────────────
    # Composants (même style que litiers)
    # ──────────────────────────────────────────

    def _accent_line(self, slide, x, y, width, color=None):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(color or PALETTE["gold"])
        line.line.fill.background()

    def _hero_metric(self, slide, x, y, w, label, value_str, sub_label=None,
                     accent_color=None, current=None, previous=None, tooltip_key=None):
        accent = accent_color or PALETTE["gold"]

        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(2.6))
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb("111111")
        bg.line.fill.background()
        bg.adjustments[0] = 0.03

        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Inches(0.25), y + Inches(0.25), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = hex_to_rgb(accent)
        dot.line.fill.background()

        self._add_textbox(
            slide, x + Inches(0.5), y + Inches(0.18), w - Inches(0.7), Inches(0.3),
            label.upper(), Pt(10), PALETTE["text_secondary"],
            font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        self._add_textbox(
            slide, x + Inches(0.25), y + Inches(0.55), w - Inches(0.5), Inches(0.8),
            value_str, Pt(44), PALETTE["white"],
            bold=True, alignment=PP_ALIGN.LEFT)

        if current is not None and previous is not None:
            try:
                curr_str = str(current).replace("%", "").replace(",", ".").strip()
                prev_str = str(previous).replace("%", "").replace(",", ".").strip()
                curr_f = float(curr_str)
                prev_f = float(prev_str)
            except (TypeError, ValueError):
                curr_f, prev_f = 0, 0

            var_pct, direction = calc_variation(curr_f, prev_f)
            if var_pct != 0:
                sign = "+" if direction == "up" else "-"
                self._add_textbox(
                    slide, x + Inches(0.25), y + Inches(1.45), w - Inches(0.5), Inches(0.3),
                    f"{sign}{var_pct}% vs mois précédent", Pt(14), accent,
                    bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)
            elif prev_f > 0:
                self._add_textbox(
                    slide, x + Inches(0.25), y + Inches(1.45), w - Inches(0.5), Inches(0.3),
                    "Stable vs mois précédent", Pt(14), accent,
                    bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        if sub_label:
            self._add_textbox(
                slide, x + Inches(0.25), y + Inches(1.85), w - Inches(0.5), Inches(0.3),
                sub_label, Pt(10), PALETTE["text_secondary"],
                italic=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)
        elif tooltip_key and tooltip_key in METRIC_TOOLTIPS:
            self._add_textbox(
                slide, x + Inches(0.25), y + Inches(1.85), w - Inches(0.5), Inches(0.6),
                METRIC_TOOLTIPS[tooltip_key], Pt(10), PALETTE["text_secondary"],
                italic=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

    def _hero_row(self, slide, metrics, y, n_cols=3):
        spacing = Inches(0.3)
        card_w = (SLIDE_WIDTH - Inches(1.2) - (n_cols - 1) * spacing) / n_cols
        start_x = (SLIDE_WIDTH - (n_cols * card_w + (n_cols - 1) * spacing)) / 2

        for i, m in enumerate(metrics):
            x = start_x + i * (card_w + spacing)
            self._hero_metric(
                slide, x, Inches(y), card_w,
                label=m.get("label", ""),
                value_str=m.get("value", "0"),
                sub_label=m.get("sub_label"),
                accent_color=m.get("accent"),
                current=m.get("current"),
                previous=m.get("previous"),
                tooltip_key=m.get("tooltip"),
            )

    def _modern_header(self, slide, title, subtitle=None, accent_color=None):
        self._set_slide_bg(slide)
        accent = accent_color or PALETTE["gold"]

        self._add_textbox(
            slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.55),
            title.upper(), Pt(26), PALETTE["white"],
            bold=True, alignment=PP_ALIGN.LEFT)

        self._accent_line(slide, Inches(0.6), Inches(0.78), Inches(1.2), accent)

        if subtitle:
            self._add_textbox(
                slide, Inches(0.6), Inches(0.88), Inches(10), Inches(0.3),
                subtitle, Pt(12), PALETTE["text_secondary"],
                font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        if LOGO_PATH.exists():
            slide.shapes.add_picture(
                str(LOGO_PATH), Inches(11.3), Inches(0.3), height=Inches(0.22))

    def _data_row_2col(self, slide, x, y, row_w, label, val_m, val_m1, bar_color):
        row_h = Inches(0.55)

        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, row_w, row_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb("111111")
        bg.line.fill.background()
        bg.adjustments[0] = 0.08

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.02), y + Inches(0.08), Inches(0.04), row_h - Inches(0.16))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(bar_color)
        bar.line.fill.background()

        self._add_textbox(slide, x + Inches(0.2), y + Inches(0.1),
            Inches(4.5), Inches(0.35),
            label, Pt(13), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        col_w = Inches(2)
        self._add_textbox(slide, x + row_w * 0.55, y + Inches(0.1),
            col_w, Inches(0.35),
            val_m1, Pt(15), PALETTE["text_secondary"],
            font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, x + row_w * 0.78, y + Inches(0.1),
            col_w, Inches(0.35),
            val_m, Pt(15), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────
    # Slide 1 — Titre
    # ──────────────────────────────────────────

    def _title_slide(self, client_name, month_fr, image_path=None):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_slide_bg(slide)

        img_left = Inches(7)
        if image_path and os.path.exists(image_path):
            pic = slide.shapes.add_picture(
                image_path, img_left, Inches(0), height=SLIDE_HEIGHT)
            natural_width = pic.width
            panel_width = SLIDE_WIDTH - img_left
            pic.width = int(panel_width)
            if natural_width > panel_width:
                crop_total = 1.0 - float(panel_width) / float(natural_width)
                pic.crop_left = crop_total / 2
                pic.crop_right = crop_total / 2

        self._accent_line(slide, Inches(0), Inches(0), img_left)

        text_w = Inches(5.5)

        self._add_textbox(
            slide, Inches(0.8), Inches(2.2), text_w, Inches(0.4),
            "RAPPORT DES CAMPAGNES", Pt(13), PALETTE["text_secondary"],
            font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        self._accent_line(slide, Inches(0.8), Inches(2.7), Inches(2))

        self._add_textbox(
            slide, Inches(0.8), Inches(3.0), text_w, Inches(1.5),
            client_name.upper(), Pt(46), PALETTE["white"],
            bold=True, alignment=PP_ALIGN.LEFT)

        self._add_textbox(
            slide, Inches(0.8), Inches(4.8), text_w, Inches(0.5),
            month_fr, Pt(20), PALETTE["gold"],
            font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        if LOGO_PATH.exists():
            slide.shapes.add_picture(
                str(LOGO_PATH), Inches(0.8), Inches(6.4), height=Inches(0.3))

        self._add_textbox(
            slide, Inches(0.8), Inches(6.85), text_w, Inches(0.3),
            "WWW.TARMAAC.IO", Pt(9), "555555",
            font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

    # ──────────────────────────────────────────
    # Slide 2 — Récap
    # ──────────────────────────────────────────

    def _slide_recap(self, gen_curr, gen_prev, g_curr, g_prev, m_curr, m_prev, month_fr):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "État Récapitulatif", month_fr)

        diff = self._safe_get(gen_curr, "Diffusion All")
        diff_p = self._safe_get(gen_prev, "Diffusion All")
        cout_m = self._safe_get(m_curr, "Cout Facebook ADS")
        cout_m_p = self._safe_get(m_prev, "Cout Facebook ADS")
        cpc_m = self._safe_get(m_curr, "CPC Meta")
        cpc_m_p = self._safe_get(m_prev, "CPC Meta")

        row1 = [
            {"label": "Diffusion totale", "value": format_number(diff),
             "current": diff, "previous": diff_p, "tooltip": "Diffusion All",
             "accent": PALETTE["gold"]},
        ]
        row2 = [
            {"label": "Coût Meta Ads", "value": format_currency(cout_m),
             "current": cout_m, "previous": cout_m_p, "tooltip": "Cout Facebook ADS",
             "accent": FUNCTIONAL_COLORS["meta_color"]},
            {"label": "CPC Meta", "value": format_currency(cpc_m),
             "current": cpc_m, "previous": cpc_m_p, "tooltip": "CPC Meta",
             "accent": FUNCTIONAL_COLORS["meta_color"]},
        ]
        self._hero_row(slide, row1, y=1.3, n_cols=1)
        self._hero_row(slide, row2, y=4.2, n_cols=2)

    # ──────────────────────────────────────────
    # Slide 3 — Meta Ads
    # ──────────────────────────────────────────

    def _slide_meta(self, m_curr, m_prev, month_fr):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Détails Meta Ads - Facebook et Instagram", month_fr,
                            accent_color=FUNCTIONAL_COLORS["meta_color"])

        blue = FUNCTIONAL_COLORS["meta_color"]

        impr = self._safe_get(m_curr, "Impressions Meta")
        impr_p = self._safe_get(m_prev, "Impressions Meta")
        clics = self._safe_get(m_curr, "Clics Meta")
        clics_p = self._safe_get(m_prev, "Clics Meta")
        cout = self._safe_get(m_curr, "Cout Facebook ADS")
        cout_p = self._safe_get(m_prev, "Cout Facebook ADS")

        row1 = [
            {"label": "Impressions", "value": format_number(impr),
             "current": impr, "previous": impr_p, "tooltip": "Impressions Meta", "accent": blue},
            {"label": "Total des clics", "value": format_number(clics),
             "current": clics, "previous": clics_p, "tooltip": "Clics Meta", "accent": blue},
            {"label": "Coût Meta Ads", "value": format_currency(cout),
             "current": cout, "previous": cout_p, "tooltip": "Cout Facebook ADS", "accent": blue},
        ]

        cpc = self._safe_get(m_curr, "CPC Meta")
        cpc_p = self._safe_get(m_prev, "CPC Meta")
        ctr = self._safe_get(m_curr, "CTR Meta")
        ctr_p = self._safe_get(m_prev, "CTR Meta")

        row2 = [
            {"label": "CPC Moyen", "value": format_currency(cpc),
             "current": cpc, "previous": cpc_p, "tooltip": "CPC Meta", "accent": blue},
            {"label": "CTR", "value": format_percentage(ctr),
             "current": ctr, "previous": ctr_p, "tooltip": "CTR Meta", "accent": blue},
        ]

        self._hero_row(slide, row1, y=1.3, n_cols=3)
        self._hero_row(slide, row2, y=4.2, n_cols=2)

    # ──────────────────────────────────────────
    # Évolution (4 charts)
    # ──────────────────────────────────────────

    def _slide_evo_quad(self, history, keys_titles, platform):
        slide = self.prs.slides.add_slide(self.blank_layout)
        accent = FUNCTIONAL_COLORS["meta_color"]
        self._modern_header(slide, f"{platform} — Évolution", accent_color=accent)

        positions = [
            (Inches(0.3), Inches(1.2)),
            (Inches(6.8), Inches(1.2)),
            (Inches(0.3), Inches(4.2)),
            (Inches(6.8), Inches(4.2)),
        ]
        chart_w = Inches(6)
        chart_h = Inches(2.8)

        for i, (key, title) in enumerate(keys_titles):
            if i >= 4:
                break
            has_data = any(self._safe_get(h, key) > 0 for h in history)
            if has_data:
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    charts.create_plotly_evolution_chart(history, key, title, tmp.name, line_color=accent)
                    x, y = positions[i]
                    self.add_chart_image(slide, tmp.name, x, y, chart_w, chart_h)

    # ──────────────────────────────────────────
    # Slide 4 — Interactions par campagne (placeholders)
    # ──────────────────────────────────────────

    def _slide_interactions(self, month_fr, prev_month_fr):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Détails Meta Ads - Facebook et Instagram — Interactions", month_fr,
                            accent_color=FUNCTIONAL_COLORS["meta_color"])

        row_w = Inches(12)
        start_x = (SLIDE_WIDTH - row_w) / 2
        start_y = Inches(1.4)

        col_w = Inches(2)
        self._add_textbox(slide, start_x + row_w * 0.55, start_y, col_w, Inches(0.3),
            prev_month_fr, Pt(10), PALETTE["text_secondary"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, start_x + row_w * 0.78, start_y, col_w, Inches(0.3),
            month_fr, Pt(10), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)

        blue = FUNCTIONAL_COLORS["meta_color"]
        y = start_y + Inches(0.4)

        # Titre Interest
        self._add_textbox(slide, start_x, y, Inches(4), Inches(0.35),
            "INTEREST", Pt(12), FUNCTIONAL_COLORS["meta_color"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)
        y += Inches(0.4)

        for label in INTEREST_ROWS:
            self._data_row_2col(slide, start_x, y, row_w,
                label=label, val_m="—", val_m1="—", bar_color=blue)
            y += Inches(0.72)

        y += Inches(0.25)

        # Titre Retarget
        self._add_textbox(slide, start_x, y, Inches(4), Inches(0.35),
            "RETARGET", Pt(12), FUNCTIONAL_COLORS["meta_color"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)
        y += Inches(0.4)

        for label in RETARGET_ROWS:
            self._data_row_2col(slide, start_x, y, row_w,
                label=label, val_m="—", val_m1="—", bar_color=blue)
            y += Inches(0.72)

    # ──────────────────────────────────────────
    # Slide 5 — Synthèse
    # ──────────────────────────────────────────

    def _slide_synthese(self, m_curr, m_prev, month_fr):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Synthèse", month_fr)

        clics = self._safe_get(m_curr, "Clics Meta")
        clics_p = self._safe_get(m_prev, "Clics Meta")
        cpc = self._safe_get(m_curr, "CPC Meta")
        cpc_p = self._safe_get(m_prev, "CPC Meta")
        cout_m = self._safe_get(m_curr, "Cout Facebook ADS")
        cout_m_p = self._safe_get(m_prev, "Cout Facebook ADS")

        row1 = [
            {"label": "Total Interactions", "value": "—",
             "sub_label": "Nombre total d'interactions sur les publications",
             "accent": PALETTE["gold"]},
            {"label": "Total des clics", "value": format_number(clics),
             "current": clics, "previous": clics_p, "tooltip": "Clics Meta",
             "accent": FUNCTIONAL_COLORS["meta_color"]},
        ]
        row2 = [
            {"label": "CPC", "value": format_currency(cpc),
             "current": cpc, "previous": cpc_p, "tooltip": "CPC Meta",
             "accent": FUNCTIONAL_COLORS["meta_color"]},
            {"label": "Total budget Meta", "value": format_currency(cout_m),
             "current": cout_m, "previous": cout_m_p, "tooltip": "Cout Facebook ADS",
             "accent": FUNCTIONAL_COLORS["meta_color"]},
        ]
        self._hero_row(slide, row1, y=1.3, n_cols=2)
        self._hero_row(slide, row2, y=4.2, n_cols=2)

    # ──────────────────────────────────────────
    # Slide 6 — Fin
    # ──────────────────────────────────────────

    def _end_slide(self):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_slide_bg(slide)

        self._accent_line(slide, Inches(0), Inches(0), SLIDE_WIDTH)

        if LOGO_PATH.exists():
            logo_h = Inches(0.6)
            pic = slide.shapes.add_picture(
                str(LOGO_PATH), Inches(0), Inches(2.8), height=logo_h)
            pic.left = int((SLIDE_WIDTH - pic.width) / 2)

        self._accent_line(slide,
            (SLIDE_WIDTH - Inches(2)) / 2, Inches(3.7), Inches(2))

        self._add_textbox(
            slide, Inches(0), Inches(4.0), SLIDE_WIDTH, Inches(0.4),
            "WWW.TARMAAC.IO", Pt(11), "555555",
            font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
