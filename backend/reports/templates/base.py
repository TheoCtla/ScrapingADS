"""
BaseTemplate — Classe abstraite pour tous les templates PPTX Tarmaac.
Design premium dark theme : fond noir, cartes KPI modernes, accents gold.
"""

import os
import logging
from abc import ABC, abstractmethod

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from backend.reports.styles import (
    PALETTE, FUNCTIONAL_COLORS, FONT_TITLE, FONT_BODY,
    FONT_SIZES, SLIDE_WIDTH, SLIDE_HEIGHT, MARGIN,
    METRIC_TOOLTIPS, INVERSE_METRICS,
    hex_to_rgb, format_currency, format_number, format_percentage, calc_variation,
)


class BaseTemplate(ABC):
    """Classe abstraite dont tous les templates h\u00e9ritent."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.blank_layout = self.prs.slide_layouts[6]

    @abstractmethod
    def generate(self, data: dict) -> Presentation:
        """G\u00e9n\u00e8re le rapport PPTX \u00e0 partir des donn\u00e9es."""
        ...

    # ──────────────────────────────────────────
    # Helpers internes
    # ──────────────────────────────────────────

    def _set_slide_bg(self, slide, color_hex: str = "000000"):
        """Applique un fond uni \u00e0 une slide."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(color_hex)

    def _add_textbox(self, slide, x, y, w, h, text, font_size, color_hex="FFFFFF",
                     bold=False, italic=False, font_name=None, alignment=PP_ALIGN.LEFT):
        """Ajoute une textbox simple."""
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        run = p.runs[0]
        run.font.size = font_size
        run.font.color.rgb = hex_to_rgb(color_hex)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font_name or FONT_TITLE
        return txBox

    def _safe_get(self, data: dict, key: str, default=0):
        """R\u00e9cup\u00e8re une valeur d'un dict, retourne default si absente ou None."""
        val = data.get(key, default)
        if val is None:
            return default
        return val

    # ──────────────────────────────────────────
    # Slides communes
    # ──────────────────────────────────────────

    def create_title_slide(self, client_name: str, month_fr: str):
        """Slide de titre — dark, premium, branding Tarmaac."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_slide_bg(slide)

        # Bande gold en haut
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.05),
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = hex_to_rgb(PALETTE["gold"])
        line_shape.line.fill.background()

        # WWW.TARMAAC.IO
        self._add_textbox(
            slide, Inches(10), Inches(0.3), Inches(3), Inches(0.4),
            "WWW.TARMAAC.IO", Pt(10), PALETTE["text_secondary"],
            font_name=FONT_BODY, alignment=PP_ALIGN.RIGHT,
        )

        # RAPPORT DES CAMPAGNES
        self._add_textbox(
            slide, MARGIN, Inches(2.2), Inches(10), Inches(0.6),
            "RAPPORT DES CAMPAGNES", Pt(18), PALETTE["text_secondary"],
            font_name=FONT_BODY, alignment=PP_ALIGN.LEFT,
        )

        # Nom du client
        self._add_textbox(
            slide, MARGIN, Inches(2.9), Inches(12), Inches(1.2),
            client_name.upper(), Pt(48), PALETTE["white"],
            bold=True, alignment=PP_ALIGN.LEFT,
        )

        # Mois en gold
        self._add_textbox(
            slide, MARGIN, Inches(4.2), Inches(6), Inches(0.6),
            month_fr, Pt(24), PALETTE["gold"],
            font_name=FONT_BODY, alignment=PP_ALIGN.LEFT,
        )

        # Barre gold sous le mois
        line2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN, Inches(4.9), Inches(2.5), Inches(0.03),
        )
        line2.fill.solid()
        line2.fill.fore_color.rgb = hex_to_rgb(PALETTE["gold"])
        line2.line.fill.background()

        # TARMAAC bas droite
        self._add_textbox(
            slide, Inches(10), Inches(6.8), Inches(3), Inches(0.5),
            "TARMAAC", Pt(14), PALETTE["white"],
            bold=True, font_name=FONT_TITLE, alignment=PP_ALIGN.RIGHT,
        )

    def create_end_slide(self):
        """Slide de fin — branding Tarmaac."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_slide_bg(slide)

        # Bande gold
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.05),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(PALETTE["gold"])
        line.line.fill.background()

        # TARMAAC centr\u00e9
        self._add_textbox(
            slide, Inches(0), Inches(2.8), SLIDE_WIDTH, Inches(1.2),
            "TARMAAC", Pt(54), PALETTE["white"],
            bold=True, alignment=PP_ALIGN.CENTER,
        )

        # Barre gold
        center_x = (SLIDE_WIDTH - Inches(3)) / 2
        line2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            center_x, Inches(4.0), Inches(3), Inches(0.03),
        )
        line2.fill.solid()
        line2.fill.fore_color.rgb = hex_to_rgb(PALETTE["gold"])
        line2.line.fill.background()

        # WWW.TARMAAC.IO
        self._add_textbox(
            slide, Inches(0), Inches(4.4), SLIDE_WIDTH, Inches(0.5),
            "WWW.TARMAAC.IO", Pt(12), PALETTE["text_secondary"],
            font_name=FONT_BODY, alignment=PP_ALIGN.CENTER,
        )

    # ──────────────────────────────────────────
    # Composants r\u00e9utilisables
    # ──────────────────────────────────────────

    def add_section_header(self, slide, title: str, subtitle: str = None, y: float = 0.3):
        """Header de section avec titre bold et sous-titre optionnel."""
        self._set_slide_bg(slide)

        # Titre
        self._add_textbox(
            slide, MARGIN, Inches(y), Inches(12), Inches(0.6),
            title.upper(), Pt(28), PALETTE["white"],
            bold=True, alignment=PP_ALIGN.LEFT,
        )

        # Barre gold
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN, Inches(y + 0.55), Inches(1.5), Inches(0.03),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(PALETTE["gold"])
        bar.line.fill.background()

        if subtitle:
            self._add_textbox(
                slide, MARGIN, Inches(y + 0.65), Inches(10), Inches(0.35),
                subtitle, Pt(13), PALETTE["text_secondary"],
                font_name=FONT_BODY, alignment=PP_ALIGN.LEFT,
            )

    def add_kpi_card(self, slide, x, y, w, h, label: str, value_str: str,
                     previous_value=None, current_value=None,
                     tooltip_key: str = None, bar_color: str = None):
        """
        Carte KPI premium : fond gris fonc\u00e9, barre lat\u00e9rale color\u00e9e,
        valeur principale grande, variation, description.
        """
        bar_color = bar_color or PALETTE["gold"]

        # Fond carte
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(PALETTE["medium_gray"])
        card.line.color.rgb = hex_to_rgb(PALETTE["light_gray"])
        card.line.width = Pt(0.5)
        card.adjustments[0] = 0.04

        # Barre lat\u00e9rale gauche
        bar_w = Inches(0.05)
        bar_h = h - Inches(0.4)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.02), y + Inches(0.2), bar_w, bar_h,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(bar_color)
        bar.line.fill.background()

        pad_left = Inches(0.22)
        content_w = w - Inches(0.35)

        # Label
        self._add_textbox(
            slide, x + pad_left, y + Inches(0.12), content_w, Inches(0.3),
            label, Pt(13), PALETTE["text_secondary"],
            font_name=FONT_BODY, alignment=PP_ALIGN.LEFT,
        )

        # Valeur principale — align\u00e9e avec la barre
        self._add_textbox(
            slide, x + pad_left, y + Inches(0.40), content_w, Inches(0.65),
            value_str, Pt(40), PALETTE["white"],
            bold=True, alignment=PP_ALIGN.LEFT,
        )

        # Variation M vs M-1
        if previous_value is not None and current_value is not None:
            try:
                curr = float(current_value)
                prev = float(previous_value)
            except (TypeError, ValueError):
                curr, prev = 0, 0

            var_pct, direction = calc_variation(curr, prev)

            if var_pct != 0:
                sign = "+" if direction == "up" else "-"

                var_text = f"{sign}{var_pct}% vs mois pr\u00e9c\u00e9dent"
                self._add_textbox(
                    slide, x + pad_left, y + h - Inches(0.78),
                    content_w, Inches(0.30),
                    var_text, Pt(14), PALETTE["gold"],
                    bold=True, font_name=FONT_BODY, alignment=PP_ALIGN.LEFT,
                )

        # Description m\u00e9trique
        if tooltip_key and tooltip_key in METRIC_TOOLTIPS:
            tooltip_text = METRIC_TOOLTIPS[tooltip_key]
            self._add_textbox(
                slide, x + pad_left, y + h - Inches(0.45),
                content_w, Inches(0.42),
                tooltip_text, Pt(10), PALETTE["text_secondary"],
                italic=True, font_name=FONT_BODY, alignment=PP_ALIGN.LEFT,
            )

    def add_data_row(self, slide, x, y, row_w, label: str, value_m: str, value_m1: str,
                     label_color: str = None, bar_color: str = None):
        """
        Ligne de donn\u00e9es M vs M-1 — bandeau moderne avec barre lat\u00e9rale.
        """
        label_color = label_color or PALETTE["white"]
        bar_color = bar_color or PALETTE["gold"]
        row_h = Inches(0.6)

        # Fond
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, row_w, row_h,
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb(PALETTE["medium_gray"])
        bg.line.color.rgb = hex_to_rgb(PALETTE["light_gray"])
        bg.line.width = Pt(0.5)
        bg.adjustments[0] = 0.08

        # Barre lat\u00e9rale
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.02), y + Inches(0.1), Inches(0.04), row_h - Inches(0.2),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(bar_color)
        bar.line.fill.background()

        # Label
        self._add_textbox(
            slide, x + Inches(0.2), y + Inches(0.1),
            Inches(5), Inches(0.4),
            label, Pt(14), label_color,
            bold=True, font_name=FONT_BODY, alignment=PP_ALIGN.LEFT,
        )

        # Valeur M-1
        col_m1_x = x + row_w * 0.55
        self._add_textbox(
            slide, col_m1_x, y + Inches(0.1),
            Inches(2.5), Inches(0.4),
            value_m1, Pt(18), PALETTE["text_secondary"],
            font_name=FONT_BODY, alignment=PP_ALIGN.CENTER,
        )

        # Valeur M
        col_m_x = x + row_w * 0.78
        self._add_textbox(
            slide, col_m_x, y + Inches(0.1),
            Inches(2.5), Inches(0.4),
            value_m, Pt(18), PALETTE["white"],
            bold=True, font_name=FONT_BODY, alignment=PP_ALIGN.CENTER,
        )

    def add_column_headers(self, slide, x, y, row_w, label_m1: str, label_m: str):
        """En-t\u00eates de colonnes M-1 et M."""
        col_m1_x = x + row_w * 0.55
        self._add_textbox(
            slide, col_m1_x, y, Inches(2.5), Inches(0.35),
            label_m1, Pt(12), PALETTE["text_secondary"],
            bold=True, font_name=FONT_BODY, alignment=PP_ALIGN.CENTER,
        )
        col_m_x = x + row_w * 0.78
        self._add_textbox(
            slide, col_m_x, y, Inches(2.5), Inches(0.35),
            label_m, Pt(12), PALETTE["white"],
            bold=True, font_name=FONT_BODY, alignment=PP_ALIGN.CENTER,
        )

    def add_chart_image(self, slide, image_path: str, x, y, width, height):
        """Ins\u00e8re une image PNG et supprime le fichier temporaire."""
        try:
            if os.path.exists(image_path):
                slide.shapes.add_picture(image_path, x, y, width, height)
                os.remove(image_path)
            else:
                logging.warning(f"Image non trouvee : {image_path}")
        except Exception as e:
            logging.error(f"Erreur insertion image : {e}")

    # ──────────────────────────────────────────
    # Helpers de layout
    # ──────────────────────────────────────────

    def _kpi_row(self, slide, cards_data: list, y: float,
                 card_w=None, card_h=None, start_x=None, spacing=None):
        """Dispose une rang\u00e9e de cartes KPI horizontalement, centr\u00e9e."""
        card_w = card_w or Inches(2.8)
        card_h = card_h or Inches(2.2)
        spacing = spacing or Inches(0.25)
        n = len(cards_data)

        # Centrer la rang\u00e9e
        total_w = n * card_w + (n - 1) * spacing
        if start_x is not None:
            sx = start_x
        else:
            sx = (SLIDE_WIDTH - total_w) / 2

        for i, card in enumerate(cards_data):
            cx = sx + i * (card_w + spacing)
            self.add_kpi_card(
                slide, cx, Inches(y), card_w, card_h,
                label=card.get("label", ""),
                value_str=card.get("value", "0"),
                current_value=card.get("current"),
                previous_value=card.get("previous"),
                tooltip_key=card.get("tooltip"),
                bar_color=card.get("bar_color"),
            )
