"""
Template Modern — Design ultra-minimaliste et premium.
Fond noir pur, typographie grande, accents gold tr\u00e8s subtils,
KPIs en hero numbers, pas de bordures lourdes.
"""

import os
import tempfile
import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from backend.reports.templates.base import BaseTemplate
from backend.reports.styles import (
    PALETTE, FUNCTIONAL_COLORS, MARGIN, SLIDE_WIDTH, SLIDE_HEIGHT,
    METRIC_TOOLTIPS, INVERSE_METRICS,
    hex_to_rgb, format_currency, format_number, format_percentage, calc_variation,
)
from backend.reports import charts


# Pages spécifiques litiers
LITIER_PAGES = [
    "Page Généraliste",
    "Page Lit Coffre",
    "Page Lit motorisé",
    "Page Prestige",
    "Page Lit escamotable",
]

# Mapping nom d'onglet Sheet → fichier image dans assets/
# La clé est un mot-clé (lowercase) cherché dans le nom de l'onglet.
# L'ordre compte : le premier match gagne.
COVER_IMAGES = [
    ("france literie", "img_france_literie.jpg"),
    ("place de la literie", "img_place_literie.jpg"),
    ("meubles rigaud", "img_meubles_rigaud.jpg"),
    ("my salon", "img_my_salon.jpg"),
    ("tousalon", "img_tousalon.jpg"),
    ("roche bobois", "img_roche_bobois.jpg"),
]

ASSETS_DIR = Path(__file__).parent.parent / "assets"
LOGO_PATH = ASSETS_DIR / "tarmaac_logo.png"


FALLBACK_IMAGE = "img_fl_fallback.jpg"


def _resolve_cover_image(client_name: str) -> Optional[str]:
    """Résout l'image de couverture pour un client. Fallback FL si pas d'image spécifique."""
    lower_name = client_name.lower()
    matched = False
    for keyword, filename in COVER_IMAGES:
        if keyword in lower_name:
            matched = True
            path = ASSETS_DIR / filename
            if path.exists():
                return str(path)
            break

    # Fallback uniquement pour les clients France Literie
    if matched or "france literie" in lower_name or lower_name.startswith("fl"):
        fallback = ASSETS_DIR / FALLBACK_IMAGE
        if fallback.exists():
            return str(fallback)

    return None


class TemplateModern(BaseTemplate):
    """Template ultra-moderne pour les rapports clients."""

    @staticmethod
    def _has_data(section: dict) -> bool:
        """Vérifie si une section contient au moins une valeur non nulle."""
        return any(
            v not in (0, None, "", "0", "0.0")
            for k, v in section.items()
            if k not in ("month", "month_fr")
        )

    def _get_iti(self, data: dict) -> float:
        """Récupère les itinéraires quelle que soit l'orthographe du header."""
        return (self._safe_get(data, "Demandes d'itinéraires")
                or self._safe_get(data, "Demande d'itineraires")
                or self._safe_get(data, "Itinéraires"))

    def generate(self, data: dict) -> Presentation:
        google = data.get("google_ads", {})
        meta = data.get("meta_ads", {})
        microsoft = data.get("microsoft_ads", {})
        conv = data.get("conversions", {})
        general = data.get("general", {})
        history = data.get("history", [])

        g_curr = google.get("current", {})
        g_prev = google.get("previous", {})
        m_curr = meta.get("current", {})
        m_prev = meta.get("previous", {})
        ms_curr = microsoft.get("current", {})
        ms_prev = microsoft.get("previous", {})
        c_curr = conv.get("current", {})
        c_prev = conv.get("previous", {})
        gen_curr = general.get("current", {})
        gen_prev = general.get("previous", {})

        client = data.get("client", "")
        month_fr = data.get("month_fr", "")
        prev_month_fr = data.get("previous_month_fr", "")

        # Détection des sections avec données réelles
        has_google = self._safe_get(g_curr, "Cout Google ADS") > 0
        has_meta = self._safe_get(m_curr, "Cout Facebook ADS") > 0
        has_microsoft = self._has_data(ms_curr)
        has_itineraires = self._get_iti(c_curr) > 0
        has_contacts = self._safe_get(c_curr, "Contact Meta") > 0 or self._safe_get(g_curr, "Contact") > 0
        has_history = history and len(history) >= 2

        # Slide 1 - Titre
        cover_image = data.get("cover_image") or _resolve_cover_image(client)
        self._title_slide(client, month_fr, image_path=cover_image)

        # Slide 2 - Dashboard vue globale
        self._slide_dashboard(gen_curr, gen_prev, c_curr, c_prev, g_curr,
                              history, month_fr, has_itineraires, has_contacts)

        # Slide 2b - Évolution contacts & itinéraires
        if has_history and (has_itineraires or has_contacts):
            self._slide_evo_conversions(history, has_itineraires, has_contacts)

        # Google Ads (slide + évolution)
        if has_google:
            self._slide_google(g_curr, g_prev, month_fr)
            if has_history:
                self._slide_evo_quad(history, [
                    ("Cout Google ADS", "Coût"),
                    ("Total Impressions", "Impressions"),
                    ("Total Clic", "Clics totaux"),
                    ("Total CPC moyen", "CPC Moyen"),
                ], "Google Ads")

        # Analytics (pas pour My Salon ni France Literie Alès)
        is_my_salon = "my salon" in client.lower()
        is_ales = "alès" in client.lower()
        if not is_my_salon and not is_ales:
            analytics_data = data.get("analytics", {})
            self._slide_analytics(gen_curr, gen_prev, history, month_fr, prev_month_fr, analytics_data)

        # Meta Ads (slide + évolution)
        if has_meta:
            self._slide_meta(m_curr, m_prev, gen_curr, gen_prev, month_fr)
            if has_history:
                self._slide_evo_quad(history, [
                    ("Cout Facebook ADS", "Coût"),
                    ("Impressions Meta", "Impressions"),
                    ("Clics Meta", "Clics"),
                    ("CPC Meta", "CPC"),
                ], "Meta Ads - Facebook et Instagram")

        # Microsoft Ads (slide + évolution)
        if has_microsoft:
            self._slide_microsoft(ms_curr, ms_prev, month_fr)
            if has_history:
                self._slide_evo_quad(history, [
                    ("Cout Microsoft ADS", "Coût"),
                    ("Impressions Micro", "Impressions"),
                    ("Clics Micro", "Clics"),
                    ("CPC Micro", "CPC"),
                ], "Microsoft Ads")

        # Synthèse (adaptative selon les plateformes actives)
        self._slide_synthese(
            gen_curr, gen_prev, g_curr, g_prev, m_curr, m_prev, ms_curr, ms_prev,
            month_fr, has_google, has_meta, has_microsoft,
        )

        # Fin
        self._end_slide()

        return self.prs

    # ──────────────────────────────────────────
    # Composants modernes
    # ──────────────────────────────────────────

    def _accent_line(self, slide, x, y, width, color=None):
        """Fine ligne d\u00e9corative."""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(color or PALETTE["gold"])
        line.line.fill.background()

    def _hero_metric(self, slide, x, y, w, label, value_str, sub_label=None,
                     accent_color=None, current=None, previous=None, tooltip_key=None):
        """
        M\u00e9trique hero : gros chiffre, label discret au-dessus,
        accent color\u00e9 via un point, pas de cadre visible.
        """
        accent = accent_color or PALETTE["gold"]

        # Fond ultra-subtil
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(2.6))
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb("111111")
        bg.line.fill.background()
        bg.adjustments[0] = 0.03

        # Point accent en haut \u00e0 gauche
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Inches(0.25), y + Inches(0.25), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = hex_to_rgb(accent)
        dot.line.fill.background()

        # Label
        self._add_textbox(
            slide, x + Inches(0.5), y + Inches(0.18), w - Inches(0.7), Inches(0.3),
            label.upper(), Pt(10), PALETTE["text_secondary"],
            font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        # Valeur hero
        self._add_textbox(
            slide, x + Inches(0.25), y + Inches(0.55), w - Inches(0.5), Inches(0.8),
            value_str, Pt(44), PALETTE["white"],
            bold=True, alignment=PP_ALIGN.LEFT)

        # Variation
        if current is not None and previous is not None:
            try:
                curr_str = str(current).replace("%", "").replace(",", ".").strip()
                prev_str = str(previous).replace("%", "").replace(",", ".").strip()
                curr_f = float(curr_str)
                prev_f = float(prev_str)
            except (TypeError, ValueError):
                curr_f, prev_f = 0, 0

            var_pct, direction = calc_variation(curr_f, prev_f)
            if var_pct != 0:
                sign = "+" if direction == "up" else "-"
                self._add_textbox(
                    slide, x + Inches(0.25), y + Inches(1.45), w - Inches(0.5), Inches(0.3),
                    f"{sign}{var_pct}% vs mois précédent", Pt(14), accent,
                    bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)
            elif prev_f > 0:
                self._add_textbox(
                    slide, x + Inches(0.25), y + Inches(1.45), w - Inches(0.5), Inches(0.3),
                    "Stable vs mois précédent", Pt(14), accent,
                    bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        # Description
        if tooltip_key and tooltip_key in METRIC_TOOLTIPS:
            self._add_textbox(
                slide, x + Inches(0.25), y + Inches(1.85), w - Inches(0.5), Inches(0.6),
                METRIC_TOOLTIPS[tooltip_key], Pt(10), PALETTE["text_secondary"],
                italic=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        # Sub-label optionnel
        if sub_label:
            self._add_textbox(
                slide, x + Inches(0.25), y + Inches(1.85), w - Inches(0.5), Inches(0.3),
                sub_label, Pt(10), PALETTE["text_secondary"],
                italic=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

    def _add_3col_headers(self, slide, x, y, row_w, m2_label, m1_label, m_label):
        """En-t\u00eates 3 colonnes de mois."""
        col_w = Inches(2)
        col_m2_x = x + row_w * 0.42
        col_m1_x = x + row_w * 0.60
        col_m_x = x + row_w * 0.80
        self._add_textbox(slide, col_m2_x, y, col_w, Inches(0.3),
            m2_label, Pt(10), PALETTE["text_secondary"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, col_m1_x, y, col_w, Inches(0.3),
            m1_label, Pt(10), PALETTE["text_secondary"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, col_m_x, y, col_w, Inches(0.3),
            m_label, Pt(10), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)

    def _data_row_3col(self, slide, x, y, row_w, label, val_m, val_m1, val_m2, bar_color):
        """Ligne de donn\u00e9es avec 3 colonnes de valeurs."""
        row_h = Inches(0.55)

        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, row_w, row_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb("111111")
        bg.line.fill.background()
        bg.adjustments[0] = 0.08

        # Barre lat\u00e9rale
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.02), y + Inches(0.08), Inches(0.04), row_h - Inches(0.16))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(bar_color)
        bar.line.fill.background()

        # Label
        self._add_textbox(slide, x + Inches(0.2), y + Inches(0.1),
            Inches(4.5), Inches(0.35),
            label, Pt(13), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        col_w = Inches(2)
        # M-2
        self._add_textbox(slide, x + row_w * 0.42, y + Inches(0.1),
            col_w, Inches(0.35),
            val_m2, Pt(15), PALETTE["text_secondary"],
            font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
        # M-1
        self._add_textbox(slide, x + row_w * 0.60, y + Inches(0.1),
            col_w, Inches(0.35),
            val_m1, Pt(15), PALETTE["text_secondary"],
            font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
        # M
        self._add_textbox(slide, x + row_w * 0.80, y + Inches(0.1),
            col_w, Inches(0.35),
            val_m, Pt(15), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)

    def _hero_row(self, slide, metrics, y, n_cols=3):
        """Rang\u00e9e de hero metrics centr\u00e9e."""
        spacing = Inches(0.3)
        card_w = (SLIDE_WIDTH - Inches(1.2) - (n_cols - 1) * spacing) / n_cols
        start_x = (SLIDE_WIDTH - (n_cols * card_w + (n_cols - 1) * spacing)) / 2

        for i, m in enumerate(metrics):
            x = start_x + i * (card_w + spacing)
            self._hero_metric(
                slide, x, Inches(y), card_w,
                label=m.get("label", ""),
                value_str=m.get("value", "0"),
                sub_label=m.get("sub_label"),
                accent_color=m.get("accent"),
                current=m.get("current"),
                previous=m.get("previous"),
                tooltip_key=m.get("tooltip"),
            )

    def _modern_header(self, slide, title, subtitle=None, accent_color=None):
        """Header minimaliste : titre + fine ligne color\u00e9e."""
        self._set_slide_bg(slide)
        accent = accent_color or PALETTE["gold"]

        self._add_textbox(
            slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.55),
            title.upper(), Pt(26), PALETTE["white"],
            bold=True, alignment=PP_ALIGN.LEFT)

        self._accent_line(slide, Inches(0.6), Inches(0.78), Inches(1.2), accent)

        if subtitle:
            self._add_textbox(
                slide, Inches(0.6), Inches(0.88), Inches(10), Inches(0.3),
                subtitle, Pt(12), PALETTE["text_secondary"],
                font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        # Logo discret en haut à droite
        if LOGO_PATH.exists():
            slide.shapes.add_picture(
                str(LOGO_PATH), Inches(11.3), Inches(0.3), height=Inches(0.22))

    # ──────────────────────────────────────────
    # Slide 1 - Titre
    # ──────────────────────────────────────────

    def _title_slide(self, client_name, month_fr, image_path=None):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_slide_bg(slide)

        # ── Image plein hauteur à droite (cover fit centré) ──
        img_left = Inches(7)
        if image_path and os.path.exists(image_path):
            # 1. Ajouter avec height seul → width auto = ratio d'origine
            pic = slide.shapes.add_picture(
                image_path, img_left, Inches(0), height=SLIDE_HEIGHT)
            natural_width = pic.width

            # 2. Forcer la largeur au panneau droit
            panel_width = SLIDE_WIDTH - img_left
            pic.width = int(panel_width)

            # 3. Crop centré pour annuler la distorsion
            if natural_width > panel_width:
                crop_total = 1.0 - float(panel_width) / float(natural_width)
                pic.crop_left = crop_total / 2
                pic.crop_right = crop_total / 2

        # ── Bande gold en haut (zone texte uniquement) ──
        self._accent_line(slide, Inches(0), Inches(0), img_left)

        # ── Contenu texte à gauche ──
        text_w = Inches(5.5)

        self._add_textbox(
            slide, Inches(0.8), Inches(2.2), text_w, Inches(0.4),
            "RAPPORT DES CAMPAGNES", Pt(13), PALETTE["text_secondary"],
            font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        self._accent_line(slide, Inches(0.8), Inches(2.7), Inches(2))

        name_upper = client_name.upper()
        name_size = Pt(34) if len(name_upper) > 25 else Pt(46)
        month_y = Inches(5.2) if len(name_upper) > 25 else Inches(4.8)

        self._add_textbox(
            slide, Inches(0.8), Inches(3.0), text_w, Inches(2.0),
            name_upper, name_size, PALETTE["white"],
            bold=True, alignment=PP_ALIGN.LEFT)

        self._add_textbox(
            slide, Inches(0.8), month_y, text_w, Inches(0.5),
            month_fr, Pt(20), PALETTE["gold"],
            font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        # ── Branding bas gauche ──
        if LOGO_PATH.exists():
            slide.shapes.add_picture(
                str(LOGO_PATH), Inches(0.8), Inches(6.4), height=Inches(0.3))

        self._add_textbox(
            slide, Inches(0.8), Inches(6.85), text_w, Inches(0.3),
            "WWW.TARMAAC.IO", Pt(9), "555555",
            font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

    # ──────────────────────────────────────────
    # Slide 2 - Dashboard
    # ──────────────────────────────────────────

    def _slide_dashboard(self, gen_curr, gen_prev, c_curr, c_prev, g_curr,
                         history, month_fr, has_itineraires, has_contacts):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "État Récapitulatif", month_fr)

        diff = self._safe_get(gen_curr, "Diffusion All")
        diff_p = self._safe_get(gen_prev, "Diffusion All")
        clics = self._safe_get(gen_curr, "Clics All")
        clics_p = self._safe_get(gen_prev, "Clics All")

        conv = self._safe_get(gen_curr, "Cout par conversion majeure")
        conv_p = self._safe_get(gen_prev, "Cout par conversion majeure")

        # Ligne 1 : Diffusion totale + Total des clics + Coût/conversion
        row1 = [
            {"label": "Diffusion totale", "value": format_number(diff),
             "current": diff, "previous": diff_p, "tooltip": "Diffusion All",
             "accent": PALETTE["gold"]},
            {"label": "Total des clics", "value": format_number(clics),
             "current": clics, "previous": clics_p, "tooltip": "Clics All",
             "accent": PALETTE["gold"]},
            {"label": "Coût / conversion", "value": format_currency(conv),
             "current": conv, "previous": conv_p,
             "sub_label": "Coût par calcul d'itinéraires et contacts",
             "accent": PALETTE["gold"]},
        ]
        # Ligne 2 : Itinéraires + Contacts (adaptatif)
        row2 = []
        if has_itineraires:
            iti = self._get_iti(c_curr)
            iti_p = self._get_iti(c_prev)
            row2.append({"label": "Itinéraires", "value": format_number(iti),
                         "current": iti, "previous": iti_p,
                         "sub_label": "Nombre de clics trackés sur les boutons itinéraires",
                         "accent": PALETTE["gold"]})
        if has_contacts:
            contact = self._safe_get(g_curr, "Contact") + self._safe_get(c_curr, "Contact Meta")
            contact_p = self._safe_get(c_prev, "Contact Meta")
            row2.append({"label": "Contacts", "value": format_number(contact),
                         "current": contact, "previous": contact_p,
                         "sub_label": "Nombre de clics trackés sur les boutons contacts",
                         "accent": PALETTE["gold"]})

        if row2:
            self._hero_row(slide, row1, y=1.3, n_cols=3)
            self._hero_row(slide, row2, y=4.2, n_cols=len(row2))
        else:
            self._hero_row(slide, row1, y=2.5, n_cols=3)

    # ──────────────────────────────────────────
    # Slide 3 - R\u00e9partition & \u00c9volution
    # ──────────────────────────────────────────
    # (Int\u00e9gr\u00e9 dans le dashboard ou en slide s\u00e9par\u00e9e via evo_pair)

    # ──────────────────────────────────────────
    # Slide 2b - Évolution contacts & itinéraires
    # ──────────────────────────────────────────

    def _slide_evo_conversions(self, history, has_itineraires, has_contacts):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Évolution", accent_color=PALETTE["gold"])

        # Grille 2×2 : Coûts, Coût/conversion en haut — Itinéraires, Contacts en bas
        positions = [
            (Inches(0.3), Inches(1.2)),
            (Inches(6.8), Inches(1.2)),
            (Inches(0.3), Inches(4.2)),
            (Inches(6.8), Inches(4.2)),
        ]
        chart_w = Inches(6)
        chart_h = Inches(2.8)

        # Trouver la bonne clé itinéraires dans l'historique
        iti_key = None
        if has_itineraires:
            for k in ["Demandes d'itinéraires", "Demande d'itineraires", "Itinéraires"]:
                if any(self._safe_get(h, k) > 0 for h in history):
                    iti_key = k
                    break

        charts_to_draw = [
            ("COUT ALL", "Coûts"),
            ("Cout par conversion majeure", "Coût / conversion"),
        ]
        if iti_key:
            charts_to_draw.append((iti_key, "Itinéraires"))
        if has_contacts:
            charts_to_draw.append(("Contact", "Contacts"))

        for i, (key, title) in enumerate(charts_to_draw):
            if i >= 4:
                break
            has_data = any(self._safe_get(h, key) > 0 for h in history)
            if has_data:
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    charts.create_plotly_evolution_chart(
                        history, key, title, tmp.name, line_color=PALETTE["gold"])
                    x, y = positions[i]
                    self.add_chart_image(slide, tmp.name, x, y, chart_w, chart_h)

    # ──────────────────────────────────────────
    # Slide Google Ads
    # ──────────────────────────────────────────

    def _slide_google(self, g_curr, g_prev, month_fr):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Google Ads", month_fr, accent_color=PALETTE["green"])

        green = PALETTE["green"]

        impr = self._safe_get(g_curr, "Total Impressions")
        impr_p = self._safe_get(g_prev, "Total Impressions")
        search = self._safe_get(g_curr, "Clics search")
        search_p = self._safe_get(g_prev, "Clics search")
        pmax = self._safe_get(g_curr, "Clics Perf Max")
        pmax_p = self._safe_get(g_prev, "Clics Perf Max")

        cout = self._safe_get(g_curr, "Cout Google ADS")
        cout_p = self._safe_get(g_prev, "Cout Google ADS")
        cpc = self._safe_get(g_curr, "Total CPC moyen")
        cpc_p = self._safe_get(g_prev, "Total CPC moyen")

        row1 = [
            {"label": "Impressions", "value": format_number(impr),
             "current": impr, "previous": impr_p, "tooltip": "Total Impressions", "accent": green},
            {"label": "Clics Search", "value": format_number(search),
             "current": search, "previous": search_p, "tooltip": "Clics search", "accent": green},
            {"label": "Clics PerfMax", "value": format_number(pmax),
             "current": pmax, "previous": pmax_p, "tooltip": "Clics Perf Max", "accent": green},
        ]
        row2 = [
            {"label": "Co\u00fbt Google Ads", "value": format_currency(cout),
             "current": cout, "previous": cout_p, "tooltip": "Cout Google ADS", "accent": green},
            {"label": "CPC Moyen", "value": format_currency(cpc),
             "current": cpc, "previous": cpc_p, "tooltip": "Total CPC moyen", "accent": green},
        ]
        self._hero_row(slide, row1, y=1.3, n_cols=3)
        self._hero_row(slide, row2, y=4.2, n_cols=2)

    # ──────────────────────────────────────────
    # \u00c9volution (2 charts c\u00f4te \u00e0 c\u00f4te)
    # ──────────────────────────────────────────

    def _slide_evo_quad(self, history, keys_titles, platform):
        """4 charts d'\u00e9volution en grille 2x2 sur une slide."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        accent = PALETTE["green"] if "Google" in platform else FUNCTIONAL_COLORS["meta_color"] if "Meta" in platform else PALETTE["orange"] if "Microsoft" in platform else PALETTE["gold"]
        self._modern_header(slide, f"{platform} \u2014 \u00c9volution", accent_color=accent)

        positions = [
            (Inches(0.3), Inches(1.2)),   # haut gauche
            (Inches(6.8), Inches(1.2)),   # haut droite
            (Inches(0.3), Inches(4.2)),   # bas gauche
            (Inches(6.8), Inches(4.2)),   # bas droite
        ]
        chart_w = Inches(6)
        chart_h = Inches(2.8)

        for i, (key, title) in enumerate(keys_titles):
            if i >= 4:
                break
            has_data = any(self._safe_get(h, key) > 0 for h in history)
            if has_data:
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    charts.create_plotly_evolution_chart(history, key, title, tmp.name, line_color=accent)
                    x, y = positions[i]
                    self.add_chart_image(slide, tmp.name, x, y, chart_w, chart_h)

    # ──────────────────────────────────────────
    # Analytics
    # ──────────────────────────────────────────

    def _data_row_2col(self, slide, x, y, row_w, label, val_m, val_m1, bar_color):
        """Ligne data 2-colonnes (mois M-1 + mois M) — utilisée par la slide Analytics."""
        row_h = Inches(0.55)

        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, row_w, row_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb("111111")
        bg.line.fill.background()
        bg.adjustments[0] = 0.08

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.02), y + Inches(0.08), Inches(0.04), row_h - Inches(0.16))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(bar_color)
        bar.line.fill.background()

        self._add_textbox(slide, x + Inches(0.2), y + Inches(0.1),
            Inches(4.5), Inches(0.35),
            label, Pt(13), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        col_w = Inches(2)
        self._add_textbox(slide, x + row_w * 0.55, y + Inches(0.1),
            col_w, Inches(0.35),
            val_m1, Pt(15), PALETTE["text_secondary"],
            font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, x + row_w * 0.78, y + Inches(0.1),
            col_w, Inches(0.35),
            val_m, Pt(15), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)

    def _slide_analytics(self, gen_curr, gen_prev, history, month_fr, prev_month_fr, analytics_data=None):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Détails Analytics", accent_color=PALETTE["gold"])

        self._add_textbox(
            slide, Inches(0.6), Inches(0.88), Inches(10), Inches(0.3),
            "Nombre de vues par page de destination sur la période", Pt(11),
            PALETTE["text_secondary"], italic=True, font_name="Calibri Light",
            alignment=PP_ALIGN.LEFT)

        row_w = Inches(12)
        start_x = (SLIDE_WIDTH - row_w) / 2
        start_y = Inches(1.5)

        # 2 colonnes de mois (M-1 puis M)
        col_w = Inches(2)
        self._add_textbox(slide, start_x + row_w * 0.55, start_y, col_w, Inches(0.3),
            prev_month_fr, Pt(10), PALETTE["text_secondary"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, start_x + row_w * 0.78, start_y, col_w, Inches(0.3),
            month_fr, Pt(10), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)

        # Pages : data r\u00e9elle si dispo (config GA4 du client), sinon placeholders legacy
        ga_pages = (analytics_data or {}).get("pages") or []
        if ga_pages:
            page_rows = [
                (p["label"], p.get("current", 0), p.get("previous", 0))
                for p in ga_pages
            ]
        else:
            page_rows = [(p, None, None) for p in LITIER_PAGES]

        for i, (label, val_curr, val_prev) in enumerate(page_rows):
            y = start_y + Inches(0.4) + i * Inches(0.72)
            if val_curr is None:
                self._data_row_2col(slide, start_x, y, row_w,
                    label=label, val_m="\u2014", val_m1="\u2014",
                    bar_color=PALETTE["gold"])
            else:
                self._data_row_2col(slide, start_x, y, row_w,
                    label=label,
                    val_m=format_number(val_curr),
                    val_m1=format_number(val_prev),
                    bar_color=PALETTE["gold"])

        sep_y = start_y + Inches(0.4) + len(page_rows) * Inches(0.72) + Inches(0.25)

        diff = self._safe_get(gen_curr, "Diffusion All")
        diff_p = self._safe_get(gen_prev, "Diffusion All")
        clics = self._safe_get(gen_curr, "Clics All")
        clics_p = self._safe_get(gen_prev, "Clics All")

        self._data_row_2col(slide, start_x, sep_y, row_w,
            label="Diffusion totale des publicit\u00e9s",
            val_m=format_number(diff), val_m1=format_number(diff_p),
            bar_color=PALETTE["gold"])
        self._data_row_2col(slide, start_x, sep_y + Inches(0.72), row_w,
            label="Total des clics sur les publicit\u00e9s",
            val_m=format_number(clics), val_m1=format_number(clics_p),
            bar_color=PALETTE["gold"])

    # ──────────────────────────────────────────
    # Meta Ads (tout sur une slide)
    # ──────────────────────────────────────────

    def _slide_meta(self, m_curr, m_prev, gen_curr, gen_prev, month_fr):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Détails Meta Ads - Facebook et Instagram", month_fr, accent_color=FUNCTIONAL_COLORS["meta_color"])

        blue = FUNCTIONAL_COLORS["meta_color"]

        impr = self._safe_get(m_curr, "Impressions Meta")
        impr_p = self._safe_get(m_prev, "Impressions Meta")
        clics = self._safe_get(m_curr, "Clics Meta")
        clics_p = self._safe_get(m_prev, "Clics Meta")
        cout = self._safe_get(m_curr, "Cout Facebook ADS")
        cout_p = self._safe_get(m_prev, "Cout Facebook ADS")

        cpc = self._safe_get(m_curr, "CPC Meta")
        cpc_p = self._safe_get(m_prev, "CPC Meta")
        ctr = self._safe_get(m_curr, "CTR Meta")
        ctr_p = self._safe_get(m_prev, "CTR Meta")
        cpl = self._safe_get(m_curr, "CPL Meta")
        cpl_p = self._safe_get(m_prev, "CPL Meta")

        row1 = [
            {"label": "Coût", "value": format_currency(cout),
             "current": cout, "previous": cout_p, "tooltip": "Cout Facebook ADS", "accent": blue},
            {"label": "Impressions", "value": format_number(impr),
             "current": impr, "previous": impr_p, "tooltip": "Impressions Meta", "accent": blue},
            {"label": "Clics", "value": format_number(clics),
             "current": clics, "previous": clics_p, "tooltip": "Clics Meta", "accent": blue},
        ]
        row2 = [
            {"label": "CPC", "value": format_currency(cpc),
             "current": cpc, "previous": cpc_p, "tooltip": "CPC Meta", "accent": blue},
            {"label": "CTR", "value": format_percentage(ctr),
             "current": ctr, "previous": ctr_p, "tooltip": "CTR Meta", "accent": blue},
            {"label": "Coût par calcul d'itinéraires", "value": format_currency(cpl),
             "current": cpl, "previous": cpl_p,
             "sub_label": "Coût par calcul d'itinéraires Meta",
             "accent": blue},
        ]
        self._hero_row(slide, row1, y=1.3, n_cols=3)
        self._hero_row(slide, row2, y=4.2, n_cols=3)

    # ──────────────────────────────────────────
    # Microsoft (manuel)
    # ──────────────────────────────────────────

    def _slide_microsoft(self, ms_curr, ms_prev, month_fr):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Microsoft Ads", month_fr, accent_color=PALETTE["orange"])

        o = PALETTE["orange"]

        impr = self._safe_get(ms_curr, "Impressions Micro")
        impr_p = self._safe_get(ms_prev, "Impressions Micro")
        clics = self._safe_get(ms_curr, "Clics Micro")
        clics_p = self._safe_get(ms_prev, "Clics Micro")
        cout = self._safe_get(ms_curr, "Cout Microsoft ADS")
        cout_p = self._safe_get(ms_prev, "Cout Microsoft ADS")
        cpc = self._safe_get(ms_curr, "CPC Micro")
        cpc_p = self._safe_get(ms_prev, "CPC Micro")

        row1 = [
            {"label": "Impressions", "value": format_number(impr) if impr else "\u2014",
             "current": impr, "previous": impr_p,
             "sub_label": "Nombre total d'impressions sur Microsoft", "accent": o},
            {"label": "Clics", "value": format_number(clics) if clics else "\u2014",
             "current": clics, "previous": clics_p,
             "sub_label": "Nombre total de clics sur Microsoft", "accent": o},
        ]
        row2 = [
            {"label": "Coût", "value": format_currency(cout) if cout else "\u2014",
             "current": cout, "previous": cout_p,
             "sub_label": "Coût total des campagnes sur Microsoft", "accent": o},
            {"label": "CPC Moyen", "value": format_currency(cpc) if cpc else "\u2014",
             "current": cpc, "previous": cpc_p,
             "sub_label": "Coût par clic moyen sur Microsoft", "accent": o},
        ]
        self._hero_row(slide, row1, y=1.3, n_cols=2)
        self._hero_row(slide, row2, y=4.2, n_cols=2)

    # ──────────────────────────────────────────
    # Synth\u00e8se
    # ──────────────────────────────────────────

    def _slide_synthese(self, gen_curr, gen_prev, g_curr, g_prev, m_curr, m_prev,
                        ms_curr, ms_prev, month_fr, has_google, has_meta, has_microsoft):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Synthèse", month_fr)

        # Ligne 1 (haut) : métriques générales
        cout_all = self._safe_get(gen_curr, "COUT ALL")
        cout_all_p = self._safe_get(gen_prev, "COUT ALL")
        conv = self._safe_get(gen_curr, "Cout par conversion majeure")
        conv_p = self._safe_get(gen_prev, "Cout par conversion majeure")

        row1 = [
            {"label": "Achat Média Total", "value": format_currency(cout_all),
             "current": cout_all, "previous": cout_all_p, "tooltip": "COUT ALL",
             "accent": PALETTE["gold"]},
        ]
        if conv:
            row1.append({"label": "Coût / conversion", "value": format_currency(conv),
                         "current": conv, "previous": conv_p,
                         "sub_label": "Coût total divisé par le nombre total de conversions (contacts + itinéraires)",
                         "accent": PALETTE["gold"]})
        self._hero_row(slide, row1, y=1.3, n_cols=len(row1))

        # Ligne 2 (bas) : plateformes actives
        row2 = []
        if has_google:
            cout_g = self._safe_get(g_curr, "Cout Google ADS")
            cout_g_p = self._safe_get(g_prev, "Cout Google ADS")
            row2.append({"label": "Google Ads", "value": format_currency(cout_g),
                         "current": cout_g, "previous": cout_g_p, "tooltip": "Cout Google ADS",
                         "accent": PALETTE["green"]})
        if has_meta:
            cout_m = self._safe_get(m_curr, "Cout Facebook ADS")
            cout_m_p = self._safe_get(m_prev, "Cout Facebook ADS")
            row2.append({"label": "Meta Ads", "value": format_currency(cout_m),
                         "current": cout_m, "previous": cout_m_p, "tooltip": "Cout Facebook ADS",
                         "accent": FUNCTIONAL_COLORS["meta_color"]})
        if has_microsoft:
            cout_ms = self._safe_get(ms_curr, "Cout Microsoft ADS")
            cout_ms_p = self._safe_get(ms_prev, "Cout Microsoft ADS")
            row2.append({"label": "Microsoft Ads", "value": format_currency(cout_ms),
                         "current": cout_ms, "previous": cout_ms_p,
                         "sub_label": "Budget total dépensé sur Microsoft Ads",
                         "accent": PALETTE["orange"]})

        if row2:
            self._hero_row(slide, row2, y=4.2, n_cols=len(row2))

    # ──────────────────────────────────────────
    # Fin
    # ──────────────────────────────────────────

    def _end_slide(self):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_slide_bg(slide)

        self._accent_line(slide, Inches(0), Inches(0), SLIDE_WIDTH)

        # Logo centré
        if LOGO_PATH.exists():
            logo_h = Inches(0.6)
            pic = slide.shapes.add_picture(
                str(LOGO_PATH), Inches(0), Inches(2.8), height=logo_h)
            pic.left = int((SLIDE_WIDTH - pic.width) / 2)

        self._accent_line(slide,
            (SLIDE_WIDTH - Inches(2)) / 2, Inches(3.7), Inches(2))

        self._add_textbox(
            slide, Inches(0), Inches(4.0), SLIDE_WIDTH, Inches(0.4),
            "WWW.TARMAAC.IO", Pt(11), "555555",
            font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
