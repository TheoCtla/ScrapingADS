"""
Template EvoPro — Même charte graphique que les autres templates.
Spécificités : pages Analytics spécifiques, formulaires par type (manuels),
Google + Meta, pas de Microsoft.
"""

import os
import tempfile
import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from backend.reports.templates.base import BaseTemplate
from backend.reports.styles import (
    PALETTE, FUNCTIONAL_COLORS, MARGIN, SLIDE_WIDTH, SLIDE_HEIGHT,
    METRIC_TOOLTIPS, INVERSE_METRICS,
    hex_to_rgb, format_currency, format_number, format_percentage, calc_variation,
)
from backend.reports import charts

ASSETS_DIR = Path(__file__).parent.parent / "assets"
LOGO_PATH = ASSETS_DIR / "tarmaac_logo.png"

COVER_IMAGE = "img_evopro.jpg"

# Pages Analytics spécifiques EvoPro (placeholders manuels)
ANALYTICS_PAGES = [
    "Page d'Accueil",
    "Page PC GAMER",
    "Page Parc Informatique",
    "Page Sauvegarde",
    "Page Dépannage",
    "Page Alarme",
    "Page Wifi",
]

# Formulaires par type (placeholders manuels)
FORM_TYPES = [
    "PC GAMER",
    "DÉPANNAGE",
    "CONTACT",
]


class TemplateEvopro(BaseTemplate):
    """Template pour EvoPro Informatique."""

    def _get_iti(self, data: dict) -> float:
        """Récupère les itinéraires quelle que soit l'orthographe du header."""
        return (self._safe_get(data, "Demandes d'itinéraires")
                or self._safe_get(data, "Demande d'itineraires")
                or self._safe_get(data, "Itinéraires"))

    @staticmethod
    def _has_data(section: dict) -> bool:
        return any(
            v not in (0, None, "", "0", "0.0")
            for k, v in section.items()
            if k not in ("month", "month_fr")
        )

    def generate(self, data: dict) -> Presentation:
        google = data.get("google_ads", {})
        meta = data.get("meta_ads", {})
        conv = data.get("conversions", {})
        general = data.get("general", {})
        history = data.get("history", [])

        g_curr = google.get("current", {})
        g_prev = google.get("previous", {})
        m_curr = meta.get("current", {})
        m_prev = meta.get("previous", {})
        c_curr = conv.get("current", {})
        c_prev = conv.get("previous", {})
        gen_curr = general.get("current", {})
        gen_prev = general.get("previous", {})

        client = data.get("client", "")
        month_fr = data.get("month_fr", "")
        prev_month_fr = data.get("previous_month_fr", "")

        has_google = self._safe_get(g_curr, "Cout Google ADS") > 0
        has_meta = self._safe_get(m_curr, "Cout Facebook ADS") > 0
        has_history = history and len(history) >= 2
        has_itineraires = self._get_iti(c_curr) > 0
        has_appels = self._safe_get(c_curr, "Appels Téléphoniques") > 0

        # Slide 1 — Titre
        cover_image = data.get("cover_image")
        if not cover_image:
            path = ASSETS_DIR / COVER_IMAGE
            if path.exists():
                cover_image = str(path)
        self._title_slide(client, month_fr, image_path=cover_image)

        # Slide 2 — Récap
        self._slide_recap(gen_curr, gen_prev, c_curr, c_prev, month_fr,
                          has_itineraires, has_appels)

        # Slide 2b — Évolution appels & itinéraires
        if has_history and (has_itineraires or has_appels):
            self._slide_evo_conversions(history, has_itineraires, has_appels)

        # Slide 3 — Analytics pages (placeholders)
        self._slide_analytics(gen_curr, gen_prev, month_fr, prev_month_fr)

        # Slide 4 — Formulaires (placeholders)
        self._slide_formulaires(month_fr, prev_month_fr)

        # Slide 5 — Google Ads
        if has_google:
            self._slide_google(g_curr, g_prev, month_fr)
            if has_history:
                self._slide_evo_quad(history, [
                    ("Cout Google ADS", "Coût"),
                    ("Total Impressions", "Impressions"),
                    ("Total Clic", "Clics totaux"),
                    ("Total CPC moyen", "CPC Moyen"),
                ], "Google Ads")

        # Slide 6 — Meta Ads
        if has_meta:
            self._slide_meta(m_curr, m_prev, month_fr)
            if has_history:
                self._slide_evo_quad(history, [
                    ("Cout Facebook ADS", "Coût"),
                    ("Impressions Meta", "Impressions"),
                    ("Clics Meta", "Clics"),
                    ("CPC Meta", "CPC"),
                ], "Meta Ads - Facebook et Instagram")

        # Slide 7 — Synthèse
        self._slide_synthese(gen_curr, gen_prev, g_curr, g_prev, m_curr, m_prev,
                             month_fr, has_google, has_meta)

        # Slide 8 — Fin
        self._end_slide()

        return self.prs

    # ──────────────────────────────────────────
    # Composants (même style que litiers)
    # ──────────────────────────────────────────

    def _accent_line(self, slide, x, y, width, color=None):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(color or PALETTE["gold"])
        line.line.fill.background()

    def _hero_metric(self, slide, x, y, w, label, value_str, sub_label=None,
                     accent_color=None, current=None, previous=None, tooltip_key=None):
        accent = accent_color or PALETTE["gold"]

        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(2.6))
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb("111111")
        bg.line.fill.background()
        bg.adjustments[0] = 0.03

        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Inches(0.25), y + Inches(0.25), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = hex_to_rgb(accent)
        dot.line.fill.background()

        self._add_textbox(
            slide, x + Inches(0.5), y + Inches(0.18), w - Inches(0.7), Inches(0.3),
            label.upper(), Pt(10), PALETTE["text_secondary"],
            font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        self._add_textbox(
            slide, x + Inches(0.25), y + Inches(0.55), w - Inches(0.5), Inches(0.8),
            value_str, Pt(44), PALETTE["white"],
            bold=True, alignment=PP_ALIGN.LEFT)

        if current is not None and previous is not None:
            try:
                curr_str = str(current).replace("%", "").replace(",", ".").strip()
                prev_str = str(previous).replace("%", "").replace(",", ".").strip()
                curr_f = float(curr_str)
                prev_f = float(prev_str)
            except (TypeError, ValueError):
                curr_f, prev_f = 0, 0

            var_pct, direction = calc_variation(curr_f, prev_f)
            if var_pct != 0:
                sign = "+" if direction == "up" else "-"
                self._add_textbox(
                    slide, x + Inches(0.25), y + Inches(1.45), w - Inches(0.5), Inches(0.3),
                    f"{sign}{var_pct}% vs mois précédent", Pt(14), accent,
                    bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)
            elif prev_f > 0:
                self._add_textbox(
                    slide, x + Inches(0.25), y + Inches(1.45), w - Inches(0.5), Inches(0.3),
                    "Stable vs mois précédent", Pt(14), accent,
                    bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        if sub_label:
            self._add_textbox(
                slide, x + Inches(0.25), y + Inches(1.85), w - Inches(0.5), Inches(0.3),
                sub_label, Pt(10), PALETTE["text_secondary"],
                italic=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)
        elif tooltip_key and tooltip_key in METRIC_TOOLTIPS:
            self._add_textbox(
                slide, x + Inches(0.25), y + Inches(1.85), w - Inches(0.5), Inches(0.6),
                METRIC_TOOLTIPS[tooltip_key], Pt(10), PALETTE["text_secondary"],
                italic=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

    def _hero_row(self, slide, metrics, y, n_cols=3):
        spacing = Inches(0.3)
        card_w = (SLIDE_WIDTH - Inches(1.2) - (n_cols - 1) * spacing) / n_cols
        start_x = (SLIDE_WIDTH - (n_cols * card_w + (n_cols - 1) * spacing)) / 2

        for i, m in enumerate(metrics):
            x = start_x + i * (card_w + spacing)
            self._hero_metric(
                slide, x, Inches(y), card_w,
                label=m.get("label", ""),
                value_str=m.get("value", "0"),
                sub_label=m.get("sub_label"),
                accent_color=m.get("accent"),
                current=m.get("current"),
                previous=m.get("previous"),
                tooltip_key=m.get("tooltip"),
            )

    def _modern_header(self, slide, title, subtitle=None, accent_color=None):
        self._set_slide_bg(slide)
        accent = accent_color or PALETTE["gold"]

        self._add_textbox(
            slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.55),
            title.upper(), Pt(26), PALETTE["white"],
            bold=True, alignment=PP_ALIGN.LEFT)

        self._accent_line(slide, Inches(0.6), Inches(0.78), Inches(1.2), accent)

        if subtitle:
            self._add_textbox(
                slide, Inches(0.6), Inches(0.88), Inches(10), Inches(0.3),
                subtitle, Pt(12), PALETTE["text_secondary"],
                font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        if LOGO_PATH.exists():
            slide.shapes.add_picture(
                str(LOGO_PATH), Inches(11.3), Inches(0.3), height=Inches(0.22))

    def _data_row_2col(self, slide, x, y, row_w, label, val_m, val_m1, bar_color):
        row_h = Inches(0.55)

        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, row_w, row_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb("111111")
        bg.line.fill.background()
        bg.adjustments[0] = 0.08

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.02), y + Inches(0.08), Inches(0.04), row_h - Inches(0.16))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(bar_color)
        bar.line.fill.background()

        self._add_textbox(slide, x + Inches(0.2), y + Inches(0.1),
            Inches(4.5), Inches(0.35),
            label, Pt(13), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        col_w = Inches(2)
        self._add_textbox(slide, x + row_w * 0.55, y + Inches(0.1),
            col_w, Inches(0.35),
            val_m1, Pt(15), PALETTE["text_secondary"],
            font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, x + row_w * 0.78, y + Inches(0.1),
            col_w, Inches(0.35),
            val_m, Pt(15), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────
    # Slide 1 — Titre
    # ──────────────────────────────────────────

    def _title_slide(self, client_name, month_fr, image_path=None):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_slide_bg(slide)

        img_left = Inches(7)
        if image_path and os.path.exists(image_path):
            pic = slide.shapes.add_picture(
                image_path, img_left, Inches(0), height=SLIDE_HEIGHT)
            natural_width = pic.width
            panel_width = SLIDE_WIDTH - img_left
            pic.width = int(panel_width)
            if natural_width > panel_width:
                crop_total = 1.0 - float(panel_width) / float(natural_width)
                pic.crop_left = crop_total / 2
                pic.crop_right = crop_total / 2

        self._accent_line(slide, Inches(0), Inches(0), img_left)

        text_w = Inches(5.5)

        self._add_textbox(
            slide, Inches(0.8), Inches(2.2), text_w, Inches(0.4),
            "RAPPORT DES CAMPAGNES", Pt(13), PALETTE["text_secondary"],
            font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        self._accent_line(slide, Inches(0.8), Inches(2.7), Inches(2))

        self._add_textbox(
            slide, Inches(0.8), Inches(3.0), text_w, Inches(1.5),
            client_name.upper(), Pt(46), PALETTE["white"],
            bold=True, alignment=PP_ALIGN.LEFT)

        self._add_textbox(
            slide, Inches(0.8), Inches(4.8), text_w, Inches(0.5),
            month_fr, Pt(20), PALETTE["gold"],
            font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

        if LOGO_PATH.exists():
            slide.shapes.add_picture(
                str(LOGO_PATH), Inches(0.8), Inches(6.4), height=Inches(0.3))

        self._add_textbox(
            slide, Inches(0.8), Inches(6.85), text_w, Inches(0.3),
            "WWW.TARMAAC.IO", Pt(9), "555555",
            font_name="Calibri Light", alignment=PP_ALIGN.LEFT)

    # ──────────────────────────────────────────
    # Slide 2 — Récap
    # ──────────────────────────────────────────

    def _slide_recap(self, gen_curr, gen_prev, c_curr, c_prev, month_fr,
                     has_itineraires, has_appels):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "État Récapitulatif", month_fr)

        diff = self._safe_get(gen_curr, "Diffusion All")
        diff_p = self._safe_get(gen_prev, "Diffusion All")
        clics = self._safe_get(gen_curr, "Clics All")
        clics_p = self._safe_get(gen_prev, "Clics All")

        row1 = [
            {"label": "Diffusion totale", "value": format_number(diff),
             "current": diff, "previous": diff_p, "tooltip": "Diffusion All",
             "accent": PALETTE["gold"]},
            {"label": "Total des clics", "value": format_number(clics),
             "current": clics, "previous": clics_p, "tooltip": "Clics All",
             "accent": PALETTE["gold"]},
        ]
        self._hero_row(slide, row1, y=1.3, n_cols=2)

        row2 = []
        if has_itineraires:
            iti = self._get_iti(c_curr)
            iti_p = self._get_iti(c_prev)
            row2.append({"label": "Itinéraires", "value": format_number(iti),
                         "current": iti, "previous": iti_p,
                         "sub_label": "Nombre de clics trackés sur les boutons itinéraires",
                         "accent": PALETTE["gold"]})
        if has_appels:
            appels = self._safe_get(c_curr, "Appels Téléphoniques")
            appels_p = self._safe_get(c_prev, "Appels Téléphoniques")
            row2.append({"label": "Appels téléphoniques", "value": format_number(appels),
                         "current": appels, "previous": appels_p,
                         "sub_label": "Nombre de clics trackés sur les boutons contacts",
                         "accent": PALETTE["gold"]})

        if row2:
            self._hero_row(slide, row2, y=4.2, n_cols=len(row2))

    # ──────────────────────────────────────────
    # Slide 2b — Évolution appels & itinéraires
    # ──────────────────────────────────────────

    def _slide_evo_conversions(self, history, has_itineraires, has_appels):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Conversions — Évolution", accent_color=PALETTE["gold"])

        # Trouver la bonne clé itinéraires dans l'historique
        iti_key = None
        if has_itineraires:
            for k in ["Demandes d'itinéraires", "Demande d'itineraires", "Itinéraires"]:
                if any(self._safe_get(h, k) > 0 for h in history):
                    iti_key = k
                    break

        charts_to_draw = []
        if iti_key:
            charts_to_draw.append((iti_key, "Itinéraires"))
        if has_appels:
            charts_to_draw.append(("Appels Téléphoniques", "Appels"))

        if len(charts_to_draw) == 2:
            positions = [
                (Inches(0.3), Inches(2.0)),
                (Inches(6.8), Inches(2.0)),
            ]
            chart_w = Inches(6)
            chart_h = Inches(4)
        else:
            positions = [
                ((SLIDE_WIDTH - Inches(8)) / 2, Inches(2.0)),
            ]
            chart_w = Inches(8)
            chart_h = Inches(4)

        for i, (key, title) in enumerate(charts_to_draw):
            has_data = any(self._safe_get(h, key) > 0 for h in history)
            if has_data:
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    charts.create_plotly_evolution_chart(
                        history, key, title, tmp.name, line_color=PALETTE["gold"])
                    x, y = positions[i]
                    self.add_chart_image(slide, tmp.name, x, y, chart_w, chart_h)

    # ──────────────────────────────────────────
    # Slide 3 — Analytics pages (placeholders)
    # ──────────────────────────────────────────

    def _slide_analytics(self, gen_curr, gen_prev, month_fr, prev_month_fr):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Répartition des pages", month_fr,
                            accent_color=PALETTE["green"])

        self._add_textbox(
            slide, Inches(0.6), Inches(0.88), Inches(10), Inches(0.3),
            "Nombre de vues par page de destination sur la période", Pt(11),
            PALETTE["text_secondary"], italic=True, font_name="Calibri Light",
            alignment=PP_ALIGN.LEFT)

        row_w = Inches(12)
        start_x = (SLIDE_WIDTH - row_w) / 2
        start_y = Inches(1.5)

        col_w = Inches(2)
        self._add_textbox(slide, start_x + row_w * 0.55, start_y, col_w, Inches(0.3),
            prev_month_fr, Pt(10), PALETTE["text_secondary"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, start_x + row_w * 0.78, start_y, col_w, Inches(0.3),
            month_fr, Pt(10), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)

        green = PALETTE["green"]
        y = start_y + Inches(0.4)

        for label in ANALYTICS_PAGES:
            self._data_row_2col(slide, start_x, y, row_w,
                label=label, val_m="—", val_m1="—", bar_color=green)
            y += Inches(0.62)

        # Totaux du Sheet
        y += Inches(0.15)
        diff = format_number(self._safe_get(gen_curr, "Diffusion All"))
        diff_p = format_number(self._safe_get(gen_prev, "Diffusion All"))
        clics = format_number(self._safe_get(gen_curr, "Clics All"))
        clics_p = format_number(self._safe_get(gen_prev, "Clics All"))

        self._data_row_2col(slide, start_x, y, row_w,
            label="Diffusion totale des publicités",
            val_m=diff, val_m1=diff_p, bar_color=PALETTE["gold"])
        y += Inches(0.62)
        self._data_row_2col(slide, start_x, y, row_w,
            label="Total des clics sur les publicités",
            val_m=clics, val_m1=clics_p, bar_color=PALETTE["gold"])

    # ──────────────────────────────────────────
    # Slide 4 — Formulaires (placeholders)
    # ──────────────────────────────────────────

    def _slide_formulaires(self, month_fr, prev_month_fr):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "État Récapitulatif — Formulaires", month_fr)

        row_w = Inches(12)
        start_x = (SLIDE_WIDTH - row_w) / 2
        start_y = Inches(1.4)

        col_w = Inches(2)
        self._add_textbox(slide, start_x + row_w * 0.55, start_y, col_w, Inches(0.3),
            prev_month_fr, Pt(10), PALETTE["text_secondary"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, start_x + row_w * 0.78, start_y, col_w, Inches(0.3),
            month_fr, Pt(10), PALETTE["white"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)

        gold = PALETTE["gold"]
        y = start_y + Inches(0.4)

        # Total formulaires (placeholder)
        self._data_row_2col(slide, start_x, y, row_w,
            label="Total formulaires", val_m="—", val_m1="—", bar_color=gold)
        y += Inches(0.9)

        # Détail par type
        self._add_textbox(slide, start_x, y, Inches(4), Inches(0.35),
            "DÉTAIL PAR TYPE", Pt(12), PALETTE["gold"],
            bold=True, font_name="Calibri Light", alignment=PP_ALIGN.LEFT)
        y += Inches(0.4)

        for label in FORM_TYPES:
            self._data_row_2col(slide, start_x, y, row_w,
                label=label, val_m="—", val_m1="—", bar_color=gold)
            y += Inches(0.72)

    # ──────────────────────────────────────────
    # Slide 5 — Google Ads
    # ──────────────────────────────────────────

    def _slide_google(self, g_curr, g_prev, month_fr):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Détails Google Ads", month_fr, accent_color=PALETTE["green"])

        green = PALETTE["green"]

        impr = self._safe_get(g_curr, "Total Impressions")
        impr_p = self._safe_get(g_prev, "Total Impressions")
        search = self._safe_get(g_curr, "Clics search")
        search_p = self._safe_get(g_prev, "Clics search")
        pmax = self._safe_get(g_curr, "Clics Perf Max")
        pmax_p = self._safe_get(g_prev, "Clics Perf Max")

        row1 = [
            {"label": "Impressions", "value": format_number(impr),
             "current": impr, "previous": impr_p, "tooltip": "Total Impressions", "accent": green},
            {"label": "Clics Search", "value": format_number(search),
             "current": search, "previous": search_p, "tooltip": "Clics search", "accent": green},
            {"label": "Clics PerfMax", "value": format_number(pmax),
             "current": pmax, "previous": pmax_p, "tooltip": "Clics Perf Max", "accent": green},
        ]

        cout = self._safe_get(g_curr, "Cout Google ADS")
        cout_p = self._safe_get(g_prev, "Cout Google ADS")
        cpc = self._safe_get(g_curr, "Total CPC moyen")
        cpc_p = self._safe_get(g_prev, "Total CPC moyen")

        row2 = [
            {"label": "Coût Google Ads", "value": format_currency(cout),
             "current": cout, "previous": cout_p, "tooltip": "Cout Google ADS", "accent": green},
            {"label": "CPC Moyen", "value": format_currency(cpc),
             "current": cpc, "previous": cpc_p, "tooltip": "Total CPC moyen", "accent": green},
        ]

        self._hero_row(slide, row1, y=1.3, n_cols=3)
        self._hero_row(slide, row2, y=4.2, n_cols=2)

    # ──────────────────────────────────────────
    # Évolution (4 charts)
    # ──────────────────────────────────────────

    def _slide_evo_quad(self, history, keys_titles, platform):
        slide = self.prs.slides.add_slide(self.blank_layout)
        if "Google" in platform:
            accent = PALETTE["green"]
        elif "Meta" in platform:
            accent = FUNCTIONAL_COLORS["meta_color"]
        else:
            accent = PALETTE["gold"]
        self._modern_header(slide, f"{platform} — Évolution", accent_color=accent)

        positions = [
            (Inches(0.3), Inches(1.2)),
            (Inches(6.8), Inches(1.2)),
            (Inches(0.3), Inches(4.2)),
            (Inches(6.8), Inches(4.2)),
        ]
        chart_w = Inches(6)
        chart_h = Inches(2.8)

        for i, (key, title) in enumerate(keys_titles):
            if i >= 4:
                break
            has_data = any(self._safe_get(h, key) > 0 for h in history)
            if has_data:
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    charts.create_plotly_evolution_chart(history, key, title, tmp.name, line_color=accent)
                    x, y = positions[i]
                    self.add_chart_image(slide, tmp.name, x, y, chart_w, chart_h)

    # ──────────────────────────────────────────
    # Slide 6 — Meta Ads
    # ──────────────────────────────────────────

    def _slide_meta(self, m_curr, m_prev, month_fr):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Détails Meta Ads - Facebook et Instagram", month_fr,
                            accent_color=FUNCTIONAL_COLORS["meta_color"])

        blue = FUNCTIONAL_COLORS["meta_color"]

        impr = self._safe_get(m_curr, "Impressions Meta")
        impr_p = self._safe_get(m_prev, "Impressions Meta")
        clics = self._safe_get(m_curr, "Clics Meta")
        clics_p = self._safe_get(m_prev, "Clics Meta")
        cout = self._safe_get(m_curr, "Cout Facebook ADS")
        cout_p = self._safe_get(m_prev, "Cout Facebook ADS")

        row1 = [
            {"label": "Impressions", "value": format_number(impr),
             "current": impr, "previous": impr_p, "tooltip": "Impressions Meta", "accent": blue},
            {"label": "Total des clics", "value": format_number(clics),
             "current": clics, "previous": clics_p, "tooltip": "Clics Meta", "accent": blue},
            {"label": "Coût Meta Ads", "value": format_currency(cout),
             "current": cout, "previous": cout_p, "tooltip": "Cout Facebook ADS", "accent": blue},
        ]

        cpc = self._safe_get(m_curr, "CPC Meta")
        cpc_p = self._safe_get(m_prev, "CPC Meta")

        row2 = [
            {"label": "CPC Moyen", "value": format_currency(cpc),
             "current": cpc, "previous": cpc_p, "tooltip": "CPC Meta", "accent": blue},
        ]

        self._hero_row(slide, row1, y=1.3, n_cols=3)
        self._hero_row(slide, row2, y=4.2, n_cols=1)

    # ──────────────────────────────────────────
    # Slide 7 — Synthèse
    # ──────────────────────────────────────────

    def _slide_synthese(self, gen_curr, gen_prev, g_curr, g_prev, m_curr, m_prev,
                        month_fr, has_google, has_meta):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._modern_header(slide, "Synthèse", month_fr)

        cout_all = self._safe_get(gen_curr, "COUT ALL")
        cout_all_p = self._safe_get(gen_prev, "COUT ALL")
        conv = self._safe_get(gen_curr, "Cout par conversion majeure")
        conv_p = self._safe_get(gen_prev, "Cout par conversion majeure")

        # Ligne 1 : métriques générales
        row1 = [
            {"label": "Achat Média Total", "value": format_currency(cout_all),
             "current": cout_all, "previous": cout_all_p, "tooltip": "COUT ALL",
             "accent": PALETTE["gold"]},
        ]
        if conv:
            row1.append({"label": "Coût / conversion", "value": format_currency(conv),
                         "current": conv, "previous": conv_p, "tooltip": "Cout par conversion majeure",
                         "accent": PALETTE["gold"]})
        row1.append({"label": "Leads", "value": "—",
                     "sub_label": "Nombre de formulaires remplis sur le site web",
                     "accent": PALETTE["gold"]})
        self._hero_row(slide, row1, y=1.3, n_cols=len(row1))

        # Ligne 2 : plateformes
        row2 = []
        if has_google:
            cout_g = self._safe_get(g_curr, "Cout Google ADS")
            cout_g_p = self._safe_get(g_prev, "Cout Google ADS")
            row2.append({"label": "Google Ads", "value": format_currency(cout_g),
                         "current": cout_g, "previous": cout_g_p,
                         "sub_label": "Budget total dépensé sur Google Ads",
                         "accent": PALETTE["green"]})
        if has_meta:
            cout_m = self._safe_get(m_curr, "Cout Facebook ADS")
            cout_m_p = self._safe_get(m_prev, "Cout Facebook ADS")
            row2.append({"label": "Meta Ads", "value": format_currency(cout_m),
                         "current": cout_m, "previous": cout_m_p,
                         "sub_label": "Budget total dépensé sur Meta Ads",
                         "accent": FUNCTIONAL_COLORS["meta_color"]})

        if row2:
            self._hero_row(slide, row2, y=4.2, n_cols=len(row2))

    # ──────────────────────────────────────────
    # Slide 8 — Fin
    # ──────────────────────────────────────────

    def _end_slide(self):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._set_slide_bg(slide)

        self._accent_line(slide, Inches(0), Inches(0), SLIDE_WIDTH)

        if LOGO_PATH.exists():
            logo_h = Inches(0.6)
            pic = slide.shapes.add_picture(
                str(LOGO_PATH), Inches(0), Inches(2.8), height=logo_h)
            pic.left = int((SLIDE_WIDTH - pic.width) / 2)

        self._accent_line(slide,
            (SLIDE_WIDTH - Inches(2)) / 2, Inches(3.7), Inches(2))

        self._add_textbox(
            slide, Inches(0), Inches(4.0), SLIDE_WIDTH, Inches(0.4),
            "WWW.TARMAAC.IO", Pt(11), "555555",
            font_name="Calibri Light", alignment=PP_ALIGN.CENTER)
